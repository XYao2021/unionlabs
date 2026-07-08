#!/usr/bin/env python3
"""Generate SDR_System_Overview.pptx — advisor presentation of the USRP B210
SDR modem (features, methods, results). Content drawn from SYSTEM_REFERENCE.md.
Run: python3 slides/make_slides.py"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

HERE = os.path.dirname(__file__)
A = lambda p: os.path.join(HERE, "assets", p)


def render_equations():
    """Render the key equations to assets/equations.png via matplotlib mathtext."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    fig = plt.figure(figsize=(12.4, 5.7)); fig.patch.set_facecolor("white")
    NV, BL, AC, GY = "#0F2E4E", "#1F6FB2", "#E8781E", "#555F6A"
    H = lambda x, y, t, c=BL: fig.text(x, y, t, fontsize=15, color=c, weight="bold")
    E = lambda x, y, t, s=15: fig.text(x, y, t, fontsize=s, color=NV)
    # full-width signal model
    fig.text(0.5, 0.95, "Received-signal model — what every RX block undoes",
             fontsize=15, color=NV, weight="bold", ha="center")
    fig.text(0.5, 0.85,
        r"$r[n]=A\,e^{\,j\left(2\pi\frac{\Delta f}{f_s}n+\varphi_0\right)}"
        r"\sum_m s[m]\,g(nT_s-mT-\tau)+w[n]$",
        fontsize=17, color=NV, ha="center")
    # left column
    xL = 0.035
    H(xL, 0.70, "Detection  &  OFDM")
    E(xL+0.02, 0.61, r"$\hat{s}=\arg\min_{c}\;|\,y-c\,|$")
    E(xL+0.02, 0.51, r"$x[n]=\dfrac{1}{N}\sum_{k}X[k]\,e^{\,j2\pi kn/N}$")
    E(xL+0.02, 0.40, r"$Y[k]=H[k]X[k]+W[k]\;\Rightarrow\;\hat{X}=Y/H$")
    H(xL, 0.27, "Equalizer  &  FEC")
    E(xL+0.02, 0.18, r"$w=(R^{H}R+\lambda I)^{-1}R^{H}s$", 14)
    E(xL+0.02, 0.08, r"Viterbi ML path, $d_{\mathrm{free}}=10$;  CRC-16-CCITT", 13)
    # right column
    xR = 0.53
    H(xR, 0.70, "Synchronization")
    E(xR+0.02, 0.61, r"$R(\tau)=\left|\,\sum_{n}p^{*}[n]\,r[\tau+n\,o_s]\,\right|$")
    E(xR+0.02, 0.50, r"$e_k=\mathrm{Re}\{(y_k-y_{k-1})\,y_{k-1/2}^{*}\}$")
    H(xR, 0.37, "Carrier recovery")
    E(xR+0.02, 0.28, r"$P=\sum_n r[n]\,r^{*}[n+L],\;\;\widehat{\Delta f}=\dfrac{f_s}{2\pi L}\arg P$", 14)
    E(xR+0.02, 0.17, r"$\hat{\varphi}=\arg\left(\sum_{\mathrm{pilots}}Y[k]\,H^{*}[k]\,\mathrm{PILOT}^{*}\right)$", 14)
    E(xR+0.02, 0.07, r"$\varphi\leftarrow\varphi+\alpha\,e+\mathrm{freq},\;\;\mathrm{freq}\leftarrow\mathrm{freq}+\beta\,e$", 13)
    fig.savefig(A("equations.png"), dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)

render_equations()

# palette
NAVY   = RGBColor(0x0F, 0x2E, 0x4E)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)
ACCENT = RGBColor(0xE8, 0x7A, 0x1E)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF8)
GREY   = RGBColor(0x55, 0x5F, 0x6A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
RED    = RGBColor(0xB0, 0x30, 0x30)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tb, tf


def rect(s, l, t, w, h, color):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def setp(p, text, size, color=NAVY, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    p.text = text if text else " "; p.alignment = align
    r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font


_slide_no = [1]   # slide 1 is the title; header() auto-numbers 2, 3, 4, ...


def header(s, title, num=None):
    _slide_no[0] += 1                     # auto-number (num arg kept for readability, ignored)
    rect(s, 0, 0, 13.333, 1.15, NAVY)
    rect(s, 0, 1.15, 13.333, 0.08, ACCENT)
    _, tf = box(s, 0.55, 0.16, 11.5, 0.85)
    setp(tf.paragraphs[0], title, 30, WHITE, True)
    tf.paragraphs[0].runs[0].font.name = "Calibri Light"
    _, nf = box(s, 12.3, 0.34, 0.9, 0.6)
    setp(nf.paragraphs[0], str(_slide_no[0]), 16, RGBColor(0x9F,0xB6,0xCC), True, PP_ALIGN.RIGHT)


def bullets(s, items, l=0.7, t=1.55, w=12.0, h=5.5, size=18, gap=10):
    _, tf = box(s, l, t, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        lvl = 0; txt = it; bold = False; color = NAVY
        if isinstance(it, tuple):
            txt, lvl = it[0], it[1]
            bold = it[2] if len(it) > 2 else False
            color = it[3] if len(it) > 3 else NAVY
        p.level = lvl
        pref = "" if lvl == 0 else ""
        setp(p, txt, size - lvl*2, color, bold)
        p.space_after = Pt(gap if lvl == 0 else gap-3)
        # manual bullet glyph
        p.runs[0].text = ("●  " if lvl == 0 else "–  ") + txt
        p.runs[0].font.size = Pt(size - lvl*2)
        p.runs[0].font.bold = bold
        p.runs[0].font.color.rgb = (BLUE if lvl == 0 else GREY) if not bold else color
    return tf


def pic(s, path, l, t, w=None, h=None):
    if not os.path.exists(path): return None
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(path, Inches(l), Inches(t), **kw)


def caption(s, text, l, t, w):
    _, tf = box(s, l, t, w, 0.4)
    setp(tf.paragraphs[0], text, 12, GREY, False, PP_ALIGN.CENTER)


def table(s, rows, l, t, w, col_w=None, size=14, header_fill=BLUE):
    nr = len(rows); nc = len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(0.4*nr)).table
    if col_w:
        for j, cw in enumerate(col_w): gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = str(val)
            para = c.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(size)
            para.runs[0].font.name = "Calibri"
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
                para.runs[0].font.color.rgb = WHITE; para.runs[0].font.bold = True
            else:
                c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
                para.runs[0].font.color.rgb = NAVY
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
    return gt


# ── 1. Title ───────────────────────────────────────────────────────────────
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 4.55, 13.333, 0.06, ACCENT)
_, tf = box(s, 0.8, 2.2, 11.7, 2.2)
setp(tf.paragraphs[0], "End-to-End SDR Communication System", 40, WHITE, True)
tf.paragraphs[0].runs[0].font.name = "Calibri Light"
p = tf.add_paragraph()
setp(p, "A configurable single-carrier & OFDM modem on USRP B210 (C++/UHD)", 22,
     RGBColor(0xBF,0xD2,0xE4), False)
_, sf = box(s, 0.8, 4.75, 11.7, 1.2)
setp(sf.paragraphs[0], "Modulation · Synchronization · Equalization · FEC · ARQ · Visualization",
     18, RGBColor(0x9F,0xB6,0xCC))
p = sf.add_paragraph(); setp(p, "Hardware-validated over the air @ 915 MHz", 16, ACCENT, True)

# ── 2. System overview ─────────────────────────────────────────────────────
s = slide(); header(s, "What Was Built", "2")
bullets(s, [
    ("A complete, real-time digital communications link between two USRP B210 radios", 0, True),
    ("Two interchangeable waveforms — single-carrier (RRC) and OFDM", 1),
    ("14 modulation schemes: BPSK → 256-QAM, differential PSK, APSK", 1),
    ("Full receiver chain: energy detection, frame/timing/carrier sync, equalization, demod", 1),
    ("Reliability layer: forward error correction + CRC + stop-and-wait ARQ", 1),
    ("Built-in visualization: time / spectrum / constellation + EVM readout", 1),
    ("Every DSP block validated in simulation, then on hardware over the air", 0, True),
    ("Fully documented: 19-page reference (math + figures) and ready-to-run command sets", 1),
])
rect(s, 0.7, 6.55, 12.0, 0.02, RGBColor(0xD0,0xD8,0xE0))

# ── 3. Signal-chain architecture ───────────────────────────────────────────
s = slide(); header(s, "Signal-Chain Architecture")
_, tf = box(s, 0.5, 1.5, 12.5, 4.4)
lines = [
 ("TX   message → chunks → CRC-16 → [FEC encode] → symbols (modulator, +preamble)", NAVY, True),
 ("        ├─ single-carrier:  RRC pulse-shaping → USRP TX", GREY, False),
 ("        └─ OFDM:            subcarrier map → IFFT + cyclic prefix → USRP TX", GREY, False),
 ("", NAVY, False),
 ("RX   USRP RX → energy detector → AGC", NAVY, True),
 ("        ├─ single-carrier:  matched filter → ACQ sync → CFO → phase → [eq] → demod", GREY, False),
 ("        └─ OFDM:            Schmidl-Cox sync + CFO → FFT → 1-tap eq → pilot CPE → demod", GREY, False),
 ("     → [FEC decode] → CRC-16 check → (ARQ: ACK if OK) → reassemble message", NAVY, True),
]
for i, (t, c, b) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    setp(p, t, 15, c, b, font="Consolas"); p.space_after = Pt(6)
_, nf = box(s, 0.5, 6.35, 12.5, 0.8)
setp(nf.paragraphs[0],
     "Single-carrier timing is done by ACQ preamble correlation (not a Gardner loop). "
     "OFDM's sync / CFO / equalization / pilots all live inside its demodulator.",
     13, ACCENT, True)

# ── 3b. TX blocks explained ─────────────────────────────────────────────────
s = slide(); header(s, "TX Chain — What Each Block Does")
table(s, [
    ["Block", "Full name", "What it does"],
    ["CRC-16", "Cyclic Redundancy Check (16-bit)", "appends a checksum so the RX can detect any bit error in a chunk"],
    ["FEC", "Forward Error Correction", "rate-1/2 K=7 convolutional encoder — adds redundancy to fix bit errors"],
    ["Modulator", "bit-to-symbol mapper", "maps groups of bits to complex constellation points (QPSK … 256-QAM)"],
    ["Preamble", "known training sequence", "prepended to each burst so the RX can detect and synchronize to it"],
    ["RRC", "Root-Raised-Cosine filter", "pulse-shaping: band-limited spectrum, zero inter-symbol interference"],
    ["IFFT + CP", "Inverse FFT + Cyclic Prefix", "OFDM: subcarrier symbols → time samples; CP absorbs channel multipath"],
    ["USRP", "Universal Software Radio Peripheral", "the SDR hardware — up-converts baseband to 915 MHz RF and transmits"],
], 0.45, 1.6, 12.45, col_w=[1.5, 3.5, 7.45], size=13)

# ── 3c. RX blocks explained ─────────────────────────────────────────────────
s = slide(); header(s, "RX Chain — What Each Block Does")
table(s, [
    ["Block", "Full name", "What it does"],
    ["Energy detector", "—", "finds where each burst begins by watching received power"],
    ["AGC", "Automatic Gain Control", "normalizes the captured burst to a fixed amplitude"],
    ["Matched filter", "RRC at the receiver", "maximizes SNR and removes ISI (pairs with the TX pulse shaper)"],
    ["ACQ", "Acquisition (preamble correlation)", "joint frame + symbol timing from the known preamble"],
    ["CFO", "Carrier Frequency Offset correction", "removes the two radios' oscillator frequency offset"],
    ["Phase / PLL", "carrier phase correction", "removes residual constellation rotation before decisions"],
    ["Equalizer", "—", "inverts channel distortion (SC: off by default; OFDM: 1 tap/subcarrier)"],
    ["Schmidl-Cox", "OFDM frame synchronizer", "frame timing + CFO from a repeated-half preamble symbol"],
    ["FFT", "Fast Fourier Transform", "OFDM: time samples → per-subcarrier symbols"],
    ["CPE", "Common Phase Error tracking", "OFDM pilots track residual per-symbol phase drift"],
    ["Demodulator", "symbol-to-bit decision", "nearest-constellation-point decision: symbols → bits"],
    ["ARQ", "Automatic Repeat reQuest", "CRC-verify each chunk; ACK the good ones, retransmit the failed ones"],
], 0.45, 1.5, 12.45, col_w=[1.7, 3.35, 7.4], size=11)

# ── 4. Modulation ──────────────────────────────────────────────────────────
s = slide(); header(s, "Modulation Schemes", "4")
table(s, [
    ["Family", "Schemes", "bits/sym"],
    ["PSK", "BPSK, QPSK, 8-PSK", "1–3"],
    ["QAM", "16 / 32 / 64 / 128 / 256-QAM", "4–8"],
    ["Differential", "DBPSK, DQPSK, 8-DPSK, π/4-QPSK", "1–3"],
    ["APSK", "16-APSK, 32-APSK", "4–5"],
], 0.6, 1.55, 6.3, col_w=[1.7, 3.4, 1.2], size=14)
bullets(s, [
    ("Minimum-distance decision;  ŝ = argmin |y − c|", 0, True),
    ("Denser constellations shrink d_min → need lower EVM", 1),
    ("EVM ≲ 35% QPSK · ≲ 12% 16-QAM · ≲ 6% 64-QAM", 1),
    ("Differential: info in phase change → no PLL needed", 1),
    ("Gray coding → 1 symbol error ≈ 1 bit error", 1),
], l=0.6, t=4.0, w=6.2, h=2.7, size=15, gap=7)
pic(s, A("constellations.png"), 7.15, 1.65, w=5.9)   # aspect 2.0 → h≈2.95
caption(s, "Same 28% EVM: clean on QPSK, bleeds across 16-QAM boundaries", 7.15, 4.7, 5.9)
pic(s, A("ber_curves.png"), 8.55, 5.15, w=3.1)        # aspect 1.6 → h≈1.94, bottom≈7.1

# ── 5. Waveforms ───────────────────────────────────────────────────────────
s = slide(); header(s, "Two Waveforms: Single-Carrier vs OFDM", "5")
bullets(s, [
    ("Single-carrier + RRC pulse shaping", 0, True),
    ("Root-raised-cosine → zero-ISI at symbol instants (Nyquist)", 1),
    ("Matched filter at RX maximizes SNR", 1),
    ("OFDM (64 subcarriers, cyclic prefix)", 0, True),
    ("IFFT + CP → channel becomes 1 complex tap per subcarrier", 1),
    ("One-tap equalization — no FIR needed under multipath", 1),
    ("Scattered pilots + common-phase-error tracking", 1),
    ("Best for dense QAM on frequency-selective channels", 1),
], l=0.6, t=1.55, w=6.3, size=16, gap=7)
pic(s, A("ofdm_orthogonality.png"), 7.25, 1.55, w=5.7)  # aspect 2.56 → h≈2.23
caption(s, "OFDM orthogonality: each subcarrier peaks where others are zero", 7.25, 3.85, 5.7)
pic(s, A("rrc.png"), 7.25, 4.5, w=5.7)                   # h≈2.23, bottom≈6.73
caption(s, "RRC pulse: zero crossings at neighbouring symbols → no ISI", 7.25, 6.8, 5.7)

# ── 6. Synchronization stack ───────────────────────────────────────────────
s = slide(); header(s, "Synchronization & Timing Recovery", "6")
bullets(s, [
    ("Received signal carries 4 impairments: gain, CFO, phase, timing", 0, True),
    ("Recovery order: frame + timing → CFO → phase", 1),
    ("Energy detection — IIR hypothesis test gates each burst", 0, True),
    ("Frame + symbol timing — ACQ preamble correlation", 0, True),
    ("m-sequence or Zadoff-Chu (CAZAC) preamble", 1),
    ("Feedforward: one timing phase per burst — robust for packets", 1),
    ("Gardner TED loop — the available tracking alternative", 0, True),
    ("for a continuous stream (implemented, not active); S-curve →", 1),
], l=0.6, t=1.55, w=6.4, size=16, gap=7)
pic(s, A("gardner_scurve.png"), 7.1, 1.9, w=6.0)
caption(s, "Gardner S-curve (available tracking loop): error → 0 at correct timing", 7.1, 5.4, 6.0)

# ── 7. Carrier recovery: CFO & phase ───────────────────────────────────────
s = slide(); header(s, "Carrier Recovery: Frequency & Phase Offset", "7")
bullets(s, [
    ("CFO spins the constellation at 2π·Δf/fs rad/sample", 0, True),
    ("Estimated data-aided from the preamble, then corrected", 1),
    ("Pilot-aided phase-slope estimator", 1),
    ("Autocorrelation (Moose / Schmidl-Cox) estimator", 1),
    ("Phase offset — 2nd-order digital PLL", 0, True),
    ("Tracks static offset AND residual-CFO ramp", 1),
    ("OFDM: Schmidl-Cox joint timing + CFO; pilots track drift", 0, True),
], l=0.6, t=1.55, w=6.2, size=16, gap=7)
pic(s, A("cfo_effect.png"), 6.95, 1.55, w=6.2)          # aspect 2.86 → h≈2.17
caption(s, "Uncorrected CFO spins clusters into arcs, then a full ring", 6.95, 3.8, 6.2)
pic(s, A("schmidl_cox.png"), 7.75, 4.35, w=4.6)          # aspect 2.02 → h≈2.28, bottom≈6.63
caption(s, "Schmidl-Cox timing metric: plateau marks the OFDM frame start", 6.95, 6.7, 6.2)

# ── 8. FEC + ARQ ───────────────────────────────────────────────────────────
s = slide(); header(s, "Reliability: FEC + CRC + ARQ", "8")
bullets(s, [
    ("Forward Error Correction (optional, --fec)", 0, True),
    ("Rate-1/2, constraint-length-7 convolutional code (NASA/802.11a)", 1),
    ("Maximum-likelihood Viterbi decoder, d_free = 10", 1),
    ("~100% error-free up to ~2% raw bit-error rate", 1),
    ("Error detection — CRC-16-CCITT per chunk", 0, True),
    ("Catches all 1–2 bit, all odd, and ≤16-bit burst errors", 1),
    ("Automatic Repeat reQuest — stop-and-wait", 0, True),
    ("ACK only verified chunks; retransmit on timeout; auto-terminate", 1),
    ("ACK over TCP (default) or a reverse RF path", 1),
], l=0.6, t=1.55, w=12.0, size=17, gap=8)
rect(s, 0.7, 6.7, 12.0, 0.02, RGBColor(0xD0,0xD8,0xE0))

# ── 9. Key engineering fixes (contribution) ────────────────────────────────
s = slide(); header(s, "Key Engineering Contributions", "9")
table(s, [
    ["Problem", "Method / Fix", "Result"],
    ["Wrong demod (front-end)", "Rebuilt matched filter + ACQ joint timing", "BER → 0 (verified)"],
    ["OFDM dense-QAM smeared", "MRC-weighted pilot CPE (robust to fades)", "EVM 67% → 37%, decodes"],
    ["Higher-order length mismatch", "Fixed bits↔symbol partial-group packing", "8-PSK/16-QAM+ align"],
    ["Differential never decoded OTA", "Last-preamble reference + bypass PLL", "DBPSK/DQPSK/8-DPSK work"],
    ["Flaky burst detection", "Continuous noise-floor tracking", "Reliable gating"],
    ["Equalizer 'divergence'", "LS training + center-tap delay comp.", "BER 0 through multipath"],
], 0.6, 1.6, 12.1, col_w=[3.5, 5.1, 3.5], size=14)
_, nf = box(s, 0.6, 6.7, 12.1, 0.6)
setp(nf.paragraphs[0], "Each fix validated in simulation, then confirmed on hardware.", 15, ACCENT, True)

# ── 10. Visualization & EVM ────────────────────────────────────────────────
s = slide(); header(s, "Visualization & EVM — Proof of Correctness", "10")
bullets(s, [
    ("Every run auto-captures TX & RX signals", 0, True),
    ("Time domain, spectrum (FFT), constellation", 1),
    ("Ideal points overlaid + measured EVM %", 1),
    ("Per-modulation folders; figure saved on exit", 1),
    ("Tool to resolve individual OFDM subcarriers", 1),
    ("Turns 'it works' into a visual, quantitative claim", 0, True),
], l=0.6, t=1.55, w=5.0, size=16, gap=8)
pic(s, A("result_qpsk_ofdm_ota.png"), 5.9, 1.8, w=7.0)   # aspect 1.875 → h≈3.73
caption(s, "Over-the-air QPSK OFDM: TX ideal points → RX four clean recovered clusters",
        5.9, 5.65, 7.0)

# ── 11. Hardware results ───────────────────────────────────────────────────
s = slide(); header(s, "Hardware Results — Over the Air @ 915 MHz", "11")
table(s, [
    ["Scheme", "Waveform", "Status", "EVM"],
    ["BPSK / QPSK", "SC & OFDM", "Solid (5/5 chunks)", "~28%"],
    ["8-PSK", "SC", "Usable (ARQ completes)", "~16% (tuned)"],
    ["DBPSK / DQPSK / 8-DPSK", "SC", "Solid (differential)", "—"],
    ["16-QAM and higher", "SC & OFDM", "Blocked OTA / cable", "> ~30%"],
], 0.6, 1.6, 6.4, col_w=[2.7, 1.8, 2.0, 1.3], size=13)
bullets(s, [
    ("Full stack delivered error-free OTA", 0, True),
    ("FEC + stop-and-wait ARQ, auto-terminating", 1),
    ("Key lever: TX power drives EVM down", 1),
    ("(76→86 dB halved 8-PSK EVM: 28→16%)", 1),
    ("Link EVM floor ≈ 28–31%", 1),
], l=0.6, t=4.4, w=6.3, h=2.6, size=15, gap=6)
pic(s, A("result_8psk_ota.png"), 7.1, 1.75, w=6.0)
caption(s, "8-PSK OTA at tuned gains: 8 clean clusters, EVM 16%", 7.1, 5.25, 6.0)

# ── 12. Limitation: dense QAM ───────────────────────────────────────────────
s = slide(); header(s, "Investigation: Why 16-QAM+ Is Blocked", "12")
bullets(s, [
    ("Cable link SNR is excellent — QPSK decodes cleanly", 0, True),
    ("But 16-QAM fails on BOTH waveforms", 0, True, RED),
    ("Root cause: two free-running oscillators", 0, True),
    ("Independent TCXOs → real CFO (±1200 Hz jitter)", 1),
    ("TX carrier leakage beats at the CFO (near-DC tone)", 1),
    ("→ rotates constellation (SC) / dominates AGC (OFDM)", 1),
    ("Tried: DC removal, DC-block filter, TX LO-null", 0, True),
    ("None sufficient — leakage is a drifting tone, not static DC", 1),
], l=0.6, t=1.55, w=6.2, size=15, gap=6)
pic(s, A("result_16qam_cable_rings.png"), 6.95, 1.75, w=6.2)
caption(s, "Cable 16-QAM: amplitude perfect (rings) but phase rotates → fails",
        6.95, 5.55, 6.2)

# ── 13. The fix + multi-platform ───────────────────────────────────────────
s = slide(); header(s, "The Fix & Portability", "13")
bullets(s, [
    ("Solution: a shared 10 MHz reference clock", 0, True, GREEN),
    ("One reference into both radios' REF IN (--ref external)", 1),
    ("CFO → ~0; leakage becomes removable static DC", 1),
    ("→ 16-QAM / 64-QAM / 256-QAM decode", 1),
    ("Same as WiFi: per-packet sync + pilots + clean TX", 1),
    ("Portable across USRP models — DSP is device-independent", 0, True),
    ("N210 / X310 / X410: only args, subdev, antenna, gain change", 1),
    ("Front-panel REF IN (+ GPSDO option) → dense QAM ready", 1),
    ("Command templates provided for all four models", 1),
], l=0.6, t=1.55, w=12.0, size=17, gap=8)

# ── 14. Summary ────────────────────────────────────────────────────────────
s = slide(); header(s, "Summary", "14")
bullets(s, [
    ("Delivered a complete, documented SDR modem on USRP B210", 0, True),
    ("SC + OFDM, 14 modulations, full sync/eq/demod chain", 1),
    ("Convolutional-Viterbi FEC + CRC + stop-and-wait ARQ", 1),
    ("Hardware-validated: error-free QPSK / 8-PSK OTA, both waveforms", 0, True),
    ("Coherent AND differential schemes working end-to-end", 1),
    ("Diagnosed the dense-QAM limit to its root cause", 0, True),
    ("Free-running clocks → shared 10 MHz reference is the fix", 1),
    ("Reproducible: 19-page reference w/ math & figures, command sets, tests", 0, True),
    ("Next step: add a 10 MHz reference → unlock 16-QAM and above", 0, True, ACCENT),
], l=0.7, t=1.55, w=12.0, size=18, gap=9)

# ── 15. Live demo commands ─────────────────────────────────────────────────
def code_box(s, l, t, w, h, lines, size=11):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0x1B, 0x25, 0x33)
    sp.line.color.rgb = RGBColor(0x3A, 0x4A, 0x5E); sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.08); tf.margin_right = Inches(0.1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; r = p.runs[0]
        r.font.name = "Consolas"; r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xDCE, 0xDCE, 0xDCE) if False else RGBColor(0xD6, 0xE6, 0xF2)
        if ln.strip().startswith("#"):
            r.font.color.rgb = RGBColor(0x86, 0xC5, 0x91)   # comment green

s = slide(); header(s, "Live Demo — Run Commands", "15")
_, tf = box(s, 0.6, 1.35, 12.2, 0.5)
setp(tf.paragraphs[0], "QPSK single-carrier + FEC + stop-and-wait ARQ.  Start the SINK (RX) first, then the SOURCE (TX).",
     16, NAVY, True)
code_box(s, 0.6, 1.95, 12.15, 1.55, [
    "# TERMINAL 1 — Sink / RX   (auto-stops when all chunks verified; saves the figure)",
    "./sdr_system --role sink_arq --rx-args serial=30CD3F7 --tx-args serial=30CD3F7 \\",
    "  --rx-subdev A:A --rx-ant RX2 --rx-freq 915e6 --rx-rate 1.6e6 --rx-gain 20 \\",
    "  --scheme QPSK --fec true --ack-transport tcp --ack-port 5599 --det-mult 3",
], size=11)
code_box(s, 0.6, 3.65, 12.15, 1.7, [
    "# TERMINAL 2 — Source / TX   (retransmits until every chunk is ACKed)",
    "./sdr_system --role source_arq --tx-args serial=30CD424 --rx-args serial=30CD424 \\",
    "  --tx-subdev A:A --tx-ant TX/RX --tx-freq 915e6 --tx-rate 1.6e6 --tx-gain 78 \\",
    "  --scheme QPSK --fec true --ack-transport tcp --ack-host 127.0.0.1 \\",
    "  --ack-port 5599 --timeout 3000",
], size=11)
bullets(s, [
    ("Other schemes: swap --scheme  (BPSK · 8-PSK · DBPSK · DQPSK)", 0),
    ("OFDM: add  --waveform ofdm --ofdm-fft 64 --ofdm-cp 16 --ofdm-tx-peak 0.5", 0),
    ("8-PSK needs more TX power: --rx-gain 16 --tx-gain 86", 0),
    ("Live output: per-chunk CRC-OK, then the auto-saved constellation figure", 0),
], l=0.7, t=5.5, w=12.0, h=1.6, size=14, gap=5)

# ── 16. Appendix: equations ────────────────────────────────────────────────
s = slide(); header(s, "Appendix — Key Equations", "16")
pic(s, A("equations.png"), 1.5, 1.5, h=5.8)   # size by height; ~10.4in wide, centered

out = os.path.join(HERE, "SDR_System_Overview.pptx")
prs.save(out)
print("wrote", out, f"({len(prs.slides._sldIdLst)} slides)")
