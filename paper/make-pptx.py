#!/usr/bin/env python3
"""Build unionlabs-slides.pptx — the demo deck as an EDITABLE PowerPoint.

Mirrors unionlabs-slides.tex (10 slides). Text is native text boxes; figures
are embedded as 300-dpi PNGs rendered from figs/*.pdf. Re-run after changing
a figure:  make figs-png-300 && python3 make-pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

# ── the UnionLabs palette (figs/unionstyle.tex) ──
API   = RGBColor(0x6B, 0x4B, 0xA8)   # purple — structure/titles
CODE  = RGBColor(0x1F, 0x6F, 0x5C)   # green
PHY   = RGBColor(0x2E, 0x5E, 0x9E)   # blue
RADIO = RGBColor(0xB4, 0x70, 0x0F)   # orange
SHARE = RGBColor(0x8A, 0x2E, 0x4D)   # maroon
BODY  = RGBColor(0x22, 0x22, 0x22)
GREY  = RGBColor(0x66, 0x66, 0x66)

FIGS = "figs"
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def title(s, text, color=API):
    tf = textbox(s, Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.8))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = color
    return tf


def bullets(s, items, x, y, w, h, size=16, gap=8):
    """items: list of (text, kwargs) or plain strings. '· ' prefix handled by
    PowerPoint's own bullets so the user can edit naturally."""
    tf = textbox(s, x, y, w, h)
    first = True
    for it in items:
        text, opts = (it, {}) if isinstance(it, str) else it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        runs = text if isinstance(text, list) else [(text, {})]
        for rt, ro in runs:
            r = p.add_run(); r.text = rt
            r.font.size = Pt(opts.get("size", size))
            r.font.color.rgb = ro.get("color", opts.get("color", BODY))
            r.font.bold = ro.get("bold", opts.get("bold", False))
            r.font.italic = ro.get("italic", opts.get("italic", False))
            if ro.get("mono"):
                r.font.name = "Consolas"
        # native bullet glyph
        pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        bu = pPr.makeelement(qn('a:buChar'), {'char': '•'})
        pPr.append(bu)
    return tf


def picture(s, path, x, y, max_w, max_h):
    im = Image.open(path)
    ar = im.width / im.height
    w, h = max_w, Emu(int(max_w / ar))
    if h > max_h:
        h, w = max_h, Emu(int(max_h * ar))
    s.shapes.add_picture(path, x + (max_w - w) // 2, y + (max_h - h) // 2,
                         width=w, height=h)


M = lambda t: [(t, {"mono": True})]           # whole-run mono helper

# ── 1 · title ────────────────────────────────────────────────────────────────
s = slide()
tf = textbox(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.6))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "UnionLabs: new agent, features, physical layers"
r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = API
tf = textbox(s, Inches(1.0), Inches(4.1), Inches(11.3), Inches(1.6))
for line, sz, col in [("First Author · Second Author · Third Author", 18, BODY),
                      ("Institution", 16, GREY),
                      ("ACM MobiCom '26 Demo  ·  Austin, TX  ·  October 2026", 16, GREY)]:
    p = tf.paragraphs[0] if line.startswith("First") else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = line
    r.font.size = Pt(sz); r.font.color.rgb = col

# ── 2 · motivation ───────────────────────────────────────────────────────────
s = slide(); title(s, "The problem: everything around the algorithm")
bullets(s, [
  "A learning-over-radio experiment is easy to state and expensive to run — the algorithm is a page of code; what surrounds it is not.",
  "Platforms removed the hardware barrier, open stacks the waveform barrier …",
  "… but users still hand-link algorithms to the physical layer: hundreds of lines, days of debugging.",
  "The wiring lives in operators' heads and shell history — retyped per machine, and the retypings disagree.",
], Inches(0.8), Inches(1.5), Inches(11.6), Inches(4.2), size=18, gap=14)
tf = textbox(s, Inches(0.8), Inches(6.1), Inches(11.6), Inches(0.8))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "UnionLabs separates what an algorithm does from how the experiment is wired."
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = SHARE

# ── 3 · four parts ───────────────────────────────────────────────────────────
s = slide(); title(s, "Four designed parts")
bullets(s, [
  ([("The Agent", {"bold": True, "color": API}),
    (" — any machine becomes a managed edge node", {})], {}),
  ([("The shared workspace", {"bold": True, "color": SHARE}),
    (" — one account's state, across testbeds", {})], {}),
  ([("The algorithm contract", {"bold": True, "color": CODE}),
    (" — transmit() / receive(msg), nothing else", {})], {}),
  ([("The PHY codes", {"bold": True, "color": PHY}),
    (" — a modem whose every stage is selectable", {})], {}),
], Inches(0.7), Inches(1.7), Inches(6.4), Inches(4.8), size=18, gap=18)
picture(s, f"{FIGS}/fig-union-api-col-300-1.png",
        Inches(7.4), Inches(1.2), Inches(5.4), Inches(5.9))

# ── 4 · agent ────────────────────────────────────────────────────────────────
s = slide(); title(s, "The UnionLabs Agent")
bullets(s, [
  "One agent per edge node — server, workstation, or laptop.",
  "Install: prepares AWS, Docker, RKE2, Helm; registers with a unique edge ID + token.",
  "gRPC with the cloud: instructions down, status + heartbeats up.",
  "Owner uploads the topology via the visualizer; reservation allocates from it — live nodes only.",
  ([("Experiments become Kubernetes pods; devices exposed ", {}),
    ("only by the owner's permission", {"italic": True}), (".", {})], {}),
  "Remote access rides the cloud–edge path: pods never publicly exposed.",
], Inches(0.7), Inches(1.4), Inches(6.5), Inches(5.6), size=16, gap=12)
picture(s, f"{FIGS}/fig-agent-col-300-1.png",
        Inches(7.5), Inches(1.2), Inches(5.3), Inches(6.0))

# ── 5 · workspace ────────────────────────────────────────────────────────────
s = slide(); title(s, "One workspace across testbeds", SHARE)
bullets(s, [
  "Every session of an account mounts the same /workspace — even when the testbeds cannot reach each other.",
  ([("topologies/", {"mono": True}), (": authored wiring, written once.", {})], {}),
  ([("settings/", {"mono": True}),
    (": each session's own record — 30 s heartbeat, dropped when stale.", {})], {}),
  "Two kinds of state, two lifecycles: the shared file system never carries a fact that can go stale.",
], Inches(0.7), Inches(1.5), Inches(6.4), Inches(5.2), size=17, gap=14)
picture(s, f"{FIGS}/fig-shared-workspace-col-300-1.png",
        Inches(7.3), Inches(1.3), Inches(5.6), Inches(5.8))

# ── 6 · contract ─────────────────────────────────────────────────────────────
s = slide(); title(s, "The algorithm contract — run.sh, one command", CODE)
bullets(s, [
  ([("The algorithm writes ", {}), ("transmit()", {"mono": True}),
    (" and ", {}), ("receive(msg)", {"mono": True}),
    (" — and imports nothing from any radio driver.", {})], {}),
  "The union layer validates payloads, converts ndarray ↔ bytes, maps role names.",
  ([("The transport reads the topology + settings files; the node becomes ", {}),
    ("tx · rx · relay · peer", {"mono": True}), (".", {})], {}),
  ([("--channel", {"mono": True}), (" selects the PHY: ", {}),
    ("ideal · usrp · lora", {"mono": True}), (".", {})], {}),
  "Same path, both directions — change the physical layer with one flag, no code change.",
], Inches(0.7), Inches(1.4), Inches(6.5), Inches(5.6), size=16, gap=12)
picture(s, f"{FIGS}/fig-union-api-col-300-1.png",
        Inches(7.5), Inches(1.2), Inches(5.3), Inches(6.0))

# ── 7 · PHY ──────────────────────────────────────────────────────────────────
s = slide(); title(s, "The PHY codes", PHY)
picture(s, f"{FIGS}/fig-phy-pipeline-300-1.png",
        Inches(1.4), Inches(1.15), Inches(10.5), Inches(4.35))
bullets(s, [
  "Modulation, coding, synchronization, waveform, gains, carrier — all selectable per experiment.",
  "Stop-and-wait ARQ end to end; the ACK over TCP or a second RF path.",
], Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.5), size=16, gap=8)

# ── 8 · wiring ───────────────────────────────────────────────────────────────
s = slide(); title(s, "Wiring as a file")
bullets(s, [
  "A topology names each node, its role, its radio, its ports, and the links between nodes — every node launches from the same file.",
  ([("Per-link, per-direction media", {"bold": True}),
    (": one hop over the air, the next over Ethernet; the reply may return over TCP/IP.", {})], {}),
  ([("A wiring that cannot run is ", {}), ("refused when it is read", {"bold": True}),
    (", naming the node at fault — not minutes later, from the wrong layer.", {})], {}),
  "One command walks 44 flags and 82 topology settings from typed text to the object each configures.",
], Inches(0.8), Inches(1.6), Inches(11.6), Inches(4.8), size=18, gap=16)

# ── 9 · demonstration ────────────────────────────────────────────────────────
s = slide(); title(s, "Demonstration: each part, working")
bullets(s, [
  ([("1.  Onboarding (Agent)", {"bold": True}),
    (" — a machine enrolled on the spot; heartbeat appears; the visitor opens the container's remote desktop.", {})], {}),
  ([("2.  One workspace", {"bold": True}),
    (" — edit a wiring file here, see it from a session on a remote testbed.", {})], {}),
  ([("3.  One algorithm, different PHYs", {"bold": True}),
    (" — radio-free, then real radios by one flag; delivery, retransmissions, airtime on screen.", {})], {}),
  ([("4.  Rewiring, and refusal", {"bold": True}),
    (" — a multi-hop experiment changes medium by one field; an impossible wiring is refused by name.", {})], {}),
], Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.0), size=18, gap=18)

# ── 10 · takeaway ────────────────────────────────────────────────────────────
s = slide()
tf = textbox(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(4.6))
rows = [("Write the experiment once.", 30, True, BODY),
        ("Run it over any PHY, topology, and testbed.", 30, True, BODY),
        ("", 12, False, BODY),
        ("./run.sh --algo <yours> --topology <wiring> --node <me>", 18, False, BODY),
        ("", 12, False, BODY),
        ("One command · one wiring file · one contract", 18, False, SHARE),
        ("", 12, False, BODY),
        ("Come run it at the table.", 14, False, GREY)]
first = True
for text, sz, bold, col in rows:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col
    if text.startswith("./run.sh"):
        r.font.name = "Consolas"

prs.save("unionlabs-slides.pptx")
print("wrote unionlabs-slides.pptx,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
